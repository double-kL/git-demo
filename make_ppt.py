# -*- coding: utf-8 -*-
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree
import json

C_DARK   = RGBColor(0x1A, 0x1A, 0x2E)
C_MID    = RGBColor(0x16, 0x21, 0x3E)
C_ACCENT = RGBColor(0xE8, 0xC5, 0x6A)
C_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
C_LGRAY  = RGBColor(0xD4, 0xD4, 0xD4)
C_SUB    = RGBColor(0xA8, 0xC8, 0xE8)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def bg(slide, color=C_DARK):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, l, t, w, h, color):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    return shape


def add_textbox(slide, text, l, t, w, h, font_size=18, bold=False,
                color=C_WHITE, align=PP_ALIGN.LEFT):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Microsoft YaHei"
    return txb


def add_para_text(slide, paragraphs, l, t, w, h, font_size=13, line_sp=22,
                  color=C_LGRAY):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    first = True
    for para_text in paragraphs:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = PP_ALIGN.LEFT
        pPr = p._p.get_or_add_pPr()
        lnSpc = etree.SubElement(pPr, qn("a:lnSpc"))
        spcPts = etree.SubElement(lnSpc, qn("a:spcPts"))
        spcPts.set("val", str(line_sp * 100))
        run = p.add_run()
        run.text = para_text
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
        run.font.name = "Microsoft YaHei"
    return txb


def accent_line(slide, l, t, w):
    add_rect(slide, l, t, w, 0.04, C_ACCENT)


def hdr(slide, section, heading):
    bg(slide, C_MID)
    add_rect(slide, 0, 0, 13.33, 1.0, C_DARK)
    add_rect(slide, 0, 0, 0.08, 7.5, C_ACCENT)
    accent_line(slide, 0.08, 1.0, 13.25)
    add_textbox(slide, section, 0.2, 0.1, 9, 0.45, font_size=13, color=C_ACCENT)
    add_textbox(slide, heading, 0.2, 0.5, 12.5, 0.45,
                font_size=20, bold=True, color=C_WHITE)


def content_slide(section, heading, paras, fs=13):
    s = prs.slides.add_slide(BLANK)
    hdr(s, section, heading)
    add_para_text(s, paras, 0.35, 1.15, 12.6, 6.1, font_size=fs)
    return s


def chapter_slide(title, subtitle=""):
    s = prs.slides.add_slide(BLANK)
    bg(s, C_DARK)
    add_rect(s, 0, 0, 0.08, 7.5, C_ACCENT)
    add_rect(s, 0.08, 2.8, 13.25, 2.0, RGBColor(0x0D, 0x14, 0x26))
    accent_line(s, 0.08, 2.8, 13.25)
    accent_line(s, 0.08, 4.8, 13.25)
    add_textbox(s, title, 1.0, 3.0, 11, 1.0, font_size=40, bold=True,
                color=C_WHITE, align=PP_ALIGN.CENTER)
    if subtitle:
        add_textbox(s, subtitle, 1.0, 4.0, 11, 0.6, font_size=18,
                    color=C_SUB, align=PP_ALIGN.CENTER)
    return s


with open("slides_content.json", "r", encoding="utf-8") as f:
    data = json.load(f)

# ---- COVER ----
s = prs.slides.add_slide(BLANK)
bg(s, C_DARK)
add_rect(s, 0, 0, 0.08, 7.5, C_ACCENT)
add_rect(s, 10.5, 5.8, 2.83, 1.7, RGBColor(0x0D, 0x0D, 0x1A))
accent_line(s, 1.2, 3.55, 7.0)
add_textbox(s, data["cover"]["title"], 1.2, 1.2, 9, 1.2,
            font_size=54, bold=True, color=C_ACCENT, align=PP_ALIGN.LEFT)
add_textbox(s, data["cover"]["subtitle"], 1.2, 2.55, 10, 0.8,
            font_size=22, color=C_SUB, align=PP_ALIGN.LEFT)
add_textbox(s, data["cover"]["org"], 1.2, 3.75, 8, 0.55,
            font_size=18, color=C_LGRAY, align=PP_ALIGN.LEFT)
add_textbox(s, data["cover"]["year"], 1.2, 4.3, 4, 0.45,
            font_size=15, color=RGBColor(0x80, 0x80, 0x80), align=PP_ALIGN.LEFT)

# ---- PREFACE & ABSTRACT ----
for ch in data["chapters"]:
    chapter_slide(ch["chapter_title"], ch.get("subtitle", ""))
    for sl in ch["slides"]:
        content_slide(ch["section_label"], sl["heading"], sl["body"])

# ---- TOC ----
chapter_slide(data["toc"]["chapter_title"], data["toc"].get("subtitle", ""))
s = prs.slides.add_slide(BLANK)
hdr(s, "目录", "Contents")
for i, entry in enumerate(data["toc"]["entries"]):
    y = 1.25 + i * 1.45
    add_rect(s, 0.35, y, 0.6, 0.55, C_ACCENT)
    add_textbox(s, str(i + 1), 0.35, y, 0.6, 0.55, font_size=22, bold=True,
                color=C_DARK, align=PP_ALIGN.CENTER)
    add_textbox(s, entry["ch"], 1.1, y, 2.5, 0.4, font_size=13, color=C_ACCENT)
    add_textbox(s, entry["title"], 1.1, y + 0.32, 5.5, 0.45, font_size=18,
                bold=True, color=C_WHITE)
    add_textbox(s, entry["sub"], 6.8, y + 0.1, 6.2, 0.55,
                font_size=12, color=C_LGRAY)
    accent_line(s, 0.35, y + 0.9, 12.6)

# ---- MAIN CHAPTERS ----
for ch_data in data["main_chapters"]:
    chapter_slide(ch_data["chapter_title"], ch_data.get("subtitle", ""))
    for sl in ch_data["slides"]:
        if sl.get("two_col"):
            s = prs.slides.add_slide(BLANK)
            hdr(s, ch_data["section_label"], sl["heading"])
            add_rect(s, 6.7, 1.1, 0.04, 6.3, RGBColor(0x30, 0x40, 0x60))
            if sl.get("left_title"):
                add_textbox(s, sl["left_title"], 0.35, 1.15, 6.1, 0.4,
                            font_size=14, bold=True, color=C_ACCENT)
            if sl.get("right_title"):
                add_textbox(s, sl["right_title"], 6.85, 1.15, 6.1, 0.4,
                            font_size=14, bold=True, color=C_ACCENT)
            top = 1.6 if sl.get("left_title") else 1.15
            add_para_text(s, sl["left"], 0.35, top, 6.2, 5.6)
            add_para_text(s, sl["right"], 6.85, top, 6.2, 5.6)

        elif sl.get("data_cards"):
            s = prs.slides.add_slide(BLANK)
            hdr(s, ch_data["section_label"], sl["heading"])
            for i, card in enumerate(sl["cards"]):
                x = 0.4 + i * 3.2
                add_rect(s, x, 1.2, 2.8, 1.6, RGBColor(0x0D, 0x1A, 0x35))
                accent_line(s, x, 1.2, 2.8)
                add_textbox(s, card["value"], x, 1.35, 2.8, 0.8,
                            font_size=32, bold=True, color=C_ACCENT,
                            align=PP_ALIGN.CENTER)
                add_textbox(s, card["label"], x, 2.15, 2.8, 0.45,
                            font_size=14, color=C_LGRAY, align=PP_ALIGN.CENTER)
            add_para_text(s, sl["body"], 0.35, 2.95, 12.6, 4.3)

        elif sl.get("instrument_groups"):
            s = prs.slides.add_slide(BLANK)
            hdr(s, ch_data["section_label"], sl["heading"])
            inst_colors = [
                C_ACCENT,
                RGBColor(0x7E, 0xC8, 0xA0),
                RGBColor(0xA0, 0xC4, 0xFF),
            ]
            for i, grp in enumerate(sl["groups"]):
                color = inst_colors[i % len(inst_colors)]
                x = 0.35 + i * 4.3
                add_rect(s, x, 1.2, 4.0, 0.5, RGBColor(0x0D, 0x1A, 0x35))
                accent_line(s, x, 1.2, 4.0)
                add_textbox(s, grp["name"], x, 1.22, 4.0, 0.4,
                            font_size=18, bold=True, color=color,
                            align=PP_ALIGN.CENTER)
                for j, song in enumerate(grp["songs"]):
                    add_textbox(s, song, x + 0.1, 1.85 + j * 0.85, 3.8, 0.55,
                                font_size=14, color=C_LGRAY)
            add_textbox(s, sl.get("note", ""), 0.35, 6.3, 12.6, 0.9,
                        font_size=12, color=RGBColor(0xA0, 0xB8, 0xD0))

        else:
            content_slide(ch_data["section_label"], sl["heading"], sl["body"])

# ---- END SLIDE ----
s = prs.slides.add_slide(BLANK)
bg(s, C_DARK)
add_rect(s, 0, 0, 0.08, 7.5, C_ACCENT)
add_rect(s, 0.08, 3.1, 13.25, 1.3, RGBColor(0x0D, 0x14, 0x26))
accent_line(s, 0.08, 3.1, 13.25)
accent_line(s, 0.08, 4.4, 13.25)
add_textbox(s, data["end"]["title"], 0, 3.2, 13.33, 0.9,
            font_size=44, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
add_textbox(s, data["end"]["sub1"], 0, 4.55, 13.33, 0.5,
            font_size=17, color=C_SUB, align=PP_ALIGN.CENTER)
add_textbox(s, data["end"]["sub2"], 0, 5.15, 13.33, 0.5,
            font_size=14, color=RGBColor(0x70, 0x80, 0x90), align=PP_ALIGN.CENTER)

out = "d:/Users/vs code 代码/青苗计划PPT.pptx"
prs.save(out)
print("saved:", out)
