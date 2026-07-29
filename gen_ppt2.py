# -*- coding: utf-8 -*-
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── 颜色方案 ──────────────────────────────────────────────
C_DARK   = RGBColor(0x1A, 0x1A, 0x2E)   # 深蓝黑（背景/标题底色）
C_MID    = RGBColor(0x16, 0x21, 0x3E)   # 深蓝（内容页背景）
C_ACCENT = RGBColor(0xE8, 0xC5, 0x6A)   # 金色（装饰线/强调）
C_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
C_LGRAY  = RGBColor(0xD4, 0xD4, 0xD4)
C_SUB    = RGBColor(0xA8, 0xC8, 0xE8)   # 浅蓝（副标题）

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]  # 完全空白


# ── 工具函数 ──────────────────────────────────────────────

def bg(slide, color=C_DARK):
    """填充整页背景色"""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, l, t, w, h, color, alpha=None):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    return shape


def add_textbox(slide, text, l, t, w, h,
                font_size=18, bold=False, color=C_WHITE,
                align=PP_ALIGN.LEFT, wrap=True, line_spacing=None):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "微软雅黑"
    if line_spacing:
        from pptx.util import Pt as _Pt
        from pptx.oxml.ns import qn
        from lxml import etree
        pPr = p._p.get_or_add_pPr()
        lnSpc = etree.SubElement(pPr, qn('a:lnSpc'))
        spcPts = etree.SubElement(lnSpc, qn('a:spcPts'))
        spcPts.set('val', str(int(line_spacing * 100)))
    return txb


def add_para_text(slide, paragraphs, l, t, w, h,
                  font_size=14, color=C_LGRAY, line_spacing_pt=22):
    """多段落文本框，支持换行"""
    from pptx.oxml.ns import qn
    from lxml import etree

    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True

    first = True
    for para_text in paragraphs:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        pPr = p._p.get_or_add_pPr()
        lnSpc = etree.SubElement(pPr, qn('a:lnSpc'))
        spcPts = etree.SubElement(lnSpc, qn('a:spcPts'))
        spcPts.set('val', str(line_spacing_pt * 100))
        run = p.add_run()
        run.text = para_text
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
        run.font.name = "微软雅黑"
    return txb


def accent_line(slide, l, t, w, thick=0.04):
    """水平金色装饰线"""
    add_rect(slide, l, t, w, thick, C_ACCENT)


# ═══════════════════════════════════════════════════════════
# 1. 封面
# ═══════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, C_DARK)

# 左侧竖向金条
add_rect(s, 0, 0, 0.08, 7.5, C_ACCENT)
# 右下装饰块
add_rect(s, 10.5, 5.8, 2.83, 1.7, RGBColor(0x0D, 0x0D, 0x1A))
# 中央装饰线
accent_line(s, 1.2, 3.55, 7, 0.05)

add_textbox(s, "青苗计划", 1.2, 1.2, 9, 1.2,
            font_size=54, bold=True, color=C_ACCENT, align=PP_ALIGN.LEFT)
add_textbox(s, "KungfuMusic · 民族风流行音乐数字化传播实践", 1.2, 2.55, 10, 0.8,
            font_size=22, bold=False, color=C_SUB, align=PP_ALIGN.LEFT)
add_textbox(s, "中央民族大学音乐学院", 1.2, 3.75, 8, 0.55,
            font_size=18, color=C_LGRAY, align=PP_ALIGN.LEFT)
add_textbox(s, "2026", 1.2, 4.3, 4, 0.45,
            font_size=15, color=RGBColor(0x80,0x80,0x80), align=PP_ALIGN.LEFT)


# ═══════════════════════════════════════════════════════════
# 辅助：通用章节内页模板
# ═══════════════════════════════════════════════════════════

def chapter_slide(title, subtitle=""):
    """章节封面页"""
    s = prs.slides.add_slide(BLANK)
    bg(s, C_DARK)
    add_rect(s, 0, 0, 0.08, 7.5, C_ACCENT)
    add_rect(s, 0.08, 2.8, 13.25, 2.0, RGBColor(0x0D, 0x14, 0x26))
    accent_line(s, 0.08, 2.8, 13.25, 0.06)
    accent_line(s, 0.08, 4.8, 13.25, 0.06)
    add_textbox(s, title, 1.0, 3.0, 11, 1.0,
                font_size=40, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    if subtitle:
        add_textbox(s, subtitle, 1.0, 4.0, 11, 0.6,
                    font_size=18, color=C_SUB, align=PP_ALIGN.CENTER)
    return s


def content_slide(section_title, heading, body_paras,
                  font_size=13.5, line_spacing=20):
    """通用内容页：左侧章节色条 + 标题 + 正文"""
    s = prs.slides.add_slide(BLANK)
    bg(s, C_MID)
    # 顶栏
    add_rect(s, 0, 0, 13.33, 1.0, C_DARK)
    add_rect(s, 0, 0, 0.08, 7.5, C_ACCENT)
    accent_line(s, 0.08, 1.0, 13.25, 0.04)
    # 章节小标签
    add_textbox(s, section_title, 0.2, 0.1, 5, 0.45,
                font_size=13, color=C_ACCENT, align=PP_ALIGN.LEFT)
    # 页面大标题
    add_textbox(s, heading, 0.2, 0.5, 12.5, 0.45,
                font_size=20, bold=True, color=C_WHITE, align=PP_ALIGN.LEFT)
    # 正文
    add_para_text(s, body_paras, 0.35, 1.15, 12.6, 6.1,
                  font_size=font_size, line_spacing_pt=line_spacing)
    return s


def two_col_slide(section_title, heading, left_items, right_items,
                  left_title="", right_title=""):
    """两栏内容页"""
    s = prs.slides.add_slide(BLANK)
    bg(s, C_MID)
    add_rect(s, 0, 0, 13.33, 1.0, C_DARK)
    add_rect(s, 0, 0, 0.08, 7.5, C_ACCENT)
    accent_line(s, 0.08, 1.0, 13.25, 0.04)
    add_textbox(s, section_title, 0.2, 0.1, 5, 0.45,
                font_size=13, color=C_ACCENT, align=PP_ALIGN.LEFT)
    add_textbox(s, heading, 0.2, 0.5, 12.5, 0.45,
                font_size=20, bold=True, color=C_WHITE, align=PP_ALIGN.LEFT)
    # 分隔线
    add_rect(s, 6.7, 1.1, 0.04, 6.3, RGBColor(0x30,0x40,0x60))
    if left_title:
        add_textbox(s, left_title, 0.35, 1.15, 6.1, 0.4,
                    font_size=14, bold=True, color=C_ACCENT)
    if right_title:
        add_textbox(s, right_title, 6.85, 1.15, 6.1, 0.4,
                    font_size=14, bold=True, color=C_ACCENT)
    top = 1.6 if (left_title or right_title) else 1.15
    add_para_text(s, left_items,  0.35, top, 6.2, 5.7, font_size=13)
    add_para_text(s, right_items, 6.85, top, 6.2, 5.7, font_size=13)
    return s


# ═══════════════════════════════════════════════════════════
# 2. 前言
# ═══════════════════════════════════════════════════════════
chapter_slide("前  言", "项目背景与缘起")

content_slide(
    "前言",
    "项目简介",
    [
        "本项目以抖音原创账号"KungfuMusic"为核心，致力于民族风流行音乐的室内乐形式转化与传播。",
        "",
        "通过将古筝、马头琴、三弦、唢呐、竹笛等民族乐器融入流行音乐改编，以短视频、音频剪辑等轻量化媒介突破时空限制，构建"筛选—改编—传播—评估"全流程的民族音乐数字化传播模式。",
        "",
        "项目依托跨专业团队协作，将作曲、表演与计算机技术相融合，带动专业实践与社会实践的有机结合。",
        "",
        "前期账号发布视频作品已获得 17 万余次观看和精准粉丝积累，初步印证了民族风流行音乐室内乐转化的市场需求，为培养创新型人才、突破民乐传播瓶颈提供了有力支撑。",
    ]
)

content_slide(
    "前言",
    "时代背景与现实问题",
    [
        "在全球化与数字化深度交织的时代语境下，中华民族音乐正迎来传播机遇，也面临突出困境：",
        "",
        "▶  传播渠道单一、形式固化老旧",
        "▶  大众对中国民乐认知较为浅薄",
        "▶  传播呈现碎片化，网络影响力与文化内涵严重不匹配",
        "",
        "与此同时，大数据与网络平台的迅猛发展已催生出成功实践：",
        "柳青瑶用琵琶演绎跨界作品；《黑神话·悟空》《哪吒》系列电影借助创新表达实现文化破圈——充分验证了民族音乐创新传播的可行性与巨大潜力。",
        "",
        "中央民族大学汇聚五十六个民族的资源优势，依托音乐学院专业主导，探索"民族乐器+流行音乐+数字化短视频账号"的创新路径，具有鲜明的时代性与实践性。",
    ]
)


# ═══════════════════════════════════════════════════════════
# 3. 摘要
# ═══════════════════════════════════════════════════════════
chapter_slide("摘  要", "Abstract")

content_slide(
    "摘要",
    "项目摘要（中文）",
    [
        "为响应传承中华优秀传统文化、铸牢中华民族共同体意识的号召，本项目聚焦中华民族多民族音乐的数字化传播路径研究与实践，探索面向国内市场的创新推广模式。",
        "",
        "项目以 KungFuMusic 为核心 IP，整合音乐改编、数字制作、多平台运营与跨圈层传播等环节，构建了"创作+制作+传播"一体化实践体系。",
        "",
        "在指导教师的专业支持下，团队完成了多民族音乐作品的数字化改编、国内账号矩阵运营与用户反馈收集，验证了民族音乐在国内青年群体中的传播潜力与文化认同价值。",
        "",
        "同时，项目正视实践中存在的不足，总结了跨圈层传播中的优化方向，为后续持续深耕民族音乐 IP、完善传播策略提供了经验支撑。本项目的实践成果，为多民族音乐的数字化国内传播提供了可借鉴的路径参考，也为推动中华优秀传统文化在青年群体中的传播贡献了青年力量。",
    ]
)


# ═══════════════════════════════════════════════════════════
# 4. 目录
# ═══════════════════════════════════════════════════════════
chapter_slide("目  录", "Contents")

s = prs.slides.add_slide(BLANK)
bg(s, C_MID)
add_rect(s, 0, 0, 13.33, 1.0, C_DARK)
add_rect(s, 0, 0, 0.08, 7.5, C_ACCENT)
accent_line(s, 0.08, 1.0, 13.25, 0.04)
add_textbox(s, "目录", 0.2, 0.1, 5, 0.45, font_size=13, color=C_ACCENT)
add_textbox(s, "Contents", 0.2, 0.5, 12, 0.45, font_size=20, bold=True, color=C_WHITE)

chapters = [
    ("第一章", "研究情况介绍",       "研究背景 · 研究现状 · 研究内容与方法 · 研究创新点"),
    ("第二章", "账号运营与传播成效", "账号定位 · 运营过程 · 核心数据 · 受众反馈 · 问题调整"),
    ("第三章", "民族风流行音乐改编实践", "作品完成情况 · 乐器改编 · 录制过程 · 传播效果"),
    ("第四章", "项目整体总结与发展展望", "项目总结 · 不足反思 · 未来展望"),
]
for i, (ch, title, sub) in enumerate(chapters):
    y = 1.25 + i * 1.45
    add_rect(s, 0.35, y, 0.6, 0.55, C_ACCENT)
    add_textbox(s, str(i+1), 0.35, y, 0.6, 0.55,
                font_size=22, bold=True, color=C_DARK, align=PP_ALIGN.CENTER)
    add_textbox(s, ch, 1.1, y, 2.5, 0.4, font_size=13, color=C_ACCENT)
    add_textbox(s, title, 1.1, y+0.32, 5, 0.45, font_size=18, bold=True, color=C_WHITE)
    add_textbox(s, sub, 6.8, y+0.1, 6.2, 0.55, font_size=12, color=C_LGRAY)
    accent_line(s, 0.35, y + 0.9, 12.6, 0.02)


# ═══════════════════════════════════════════════════════════
# 5. 第一章
# ═══════════════════════════════════════════════════════════
chapter_slide("第一章", "研究情况介绍")

# 1-1 研究背景
content_slide(
    "第一章 · 研究情况介绍",
    "一、研究背景",
    [
        "中华民族音乐承载着深厚的民族文化基因，是传承中华优秀传统文化、铸牢中华民族共同体意识的重要载体。",
        "",
        "当前国内民族音乐传播存在突出问题：",
        "  ● 传播渠道单一，呈现形式固化老旧",
        "  ● 大众认知程度不高，内容传播碎片化明显",
        "  ● 民族音乐的网络影响力与文化内涵严重不匹配",
        "",
        "与此同时，大数据技术与短视频平台的普及为民族文化破圈传播提供了全新机遇：",
        "  ● 民乐达人借助网络平台演绎跨界作品",
        "  ● 《哪吒之魔童降世》《黑神话·悟空》《长安三万里》等 IP 深度融入传统民乐元素，实现文化破圈走红",
        "",
        "中央民族大学依托五十六个民族的优质文化资源与专业学术支撑，为本项目开展民族音乐创作与传播实践提供了良好平台。",
        "前期视频累计获得 17 万余次观看，有效验证了传播价值与市场需求。",
    ]
)

# 1-2 研究现状
content_slide(
    "第一章 · 研究情况介绍",
    "二、研究现状",
    [
        "【音乐改编创作】",
        "团队已完成 15 首以上改编作品的视频拍摄与初版剪辑，坚持"保留民族内核、适配室内乐质感、兼顾流行听感"的创作原则，选用竹笛、唢呐、马头琴等乐器进行融合编配。",
        "",
        "【数字化传播实践】",
        "抖音账号进入常态化运营，每周更新 2-3 条短视频，定期复盘播放量、完播率、互动率等数据，逐步优化传播策略。",
        "前期作品已获得 17 万余次观看、千余精准粉丝，充分验证了内容方向的可行性与市场潜力。",
        "",
        "【团队协作】",
        "跨专业团队优势得到充分发挥：",
        "  ● 作曲方向 → 曲目改编与室内乐编配",
        "  ● 表演方向 → 乐器录制与视频呈现",
        "  ● 计算机方向 → 后期数据统计",
        "已构建"曲目筛选—改编创作—拍摄制作—平台发布—数据复盘"的完整工作流程。",
    ]
)

# 1-3 研究内容与方法
content_slide(
    "第一章 · 研究情况介绍",
    "三、研究内容与方法",
    [
        "【四大核心研究内容】",
        "① 筛选适配短视频传播的流行音乐作品进行民族乐器室内乐改编，持续丰富作品储备",
        "② 推进抖音账号常态化运营，固定每周更新节奏，构建专业化、年轻化的民乐传播账号",
        "③ 基于传播数据与用户反馈开展效果复盘，优化改编风格、曲目选择与视频呈现形式",
        "④ 构建"创作—传播—分析"的完整实践闭环",
        "",
        "【研究方法】",
        "▶ 文献梳理法：观察线上热门视频，寻找热门影视 IP，借鉴优质民乐传播案例",
        "▶ 作品实践创作法：遵循民族器乐演奏规律与室内乐基本逻辑，对曲目进行改编打磨",
        "▶ 实证运营法：依托抖音平台开展常态化内容发布与账号运营",
        "▶ 数据归纳分析法：系统整理运营数据与用户反馈，梳理问题、总结经验，持续优化项目策略",
    ]
)

# 1-4 研究创新点
two_col_slide(
    "第一章 · 研究情况介绍",
    "四、研究创新点",
    left_title="内容与传播创新",
    right_title="实践与价值创新",
    left_items=[
        "【内容创作创新】",
        "突破传统民乐改编单一形式，将流行音乐旋律与民乐乐器编配深度融合，灵活运用竹笛、马头琴、唢呐多种民族乐器，实现传统民乐年轻化、通俗化的创新性转化。",
        "",
        "【传播模式创新】",
        "打破依赖线下演出、专业赛事的传统路径，以抖音短视频为轻量化载体，构建"创作—拍摄—发布—复盘"的闭环传播体系，有效拓宽民族音乐的大众触达渠道。",
    ],
    right_items=[
        "【实践路径创新】",
        "依托跨专业团队优势，打通作曲编曲、器乐表演、新媒体运营、数据统计全流程，形成"音乐创作+数字制作+平台传播"一体化实践模式，为高校音乐学生提供可参考的落地路径。",
        "",
        "【价值导向创新】",
        "以铸牢中华民族共同体意识为核心，立足多民族文化资源优势，用年轻化的传播方式讲好民族音乐故事，兼顾文化传承性、内容专业性与大众传播性。",
    ]
)


# ═══════════════════════════════════════════════════════════
# 6. 第二章
# ═══════════════════════════════════════════════════════════
chapter_slide("第二章", "账号运营与传播成效")

# 2-1 账号定位
content_slide(
    "第二章 · 账号运营与传播成效",
    "一、账号定位与运营策略设计",
    [
        "【核心 IP】GongfuMusic",
        "聚焦民族风流行音乐室内乐改编及国内新媒体传播，以传承传统民乐、创新国风演绎为宗旨。",
        "",
        "【目标受众】",
        "国内国风爱好者 · 民族乐器学习者 · 高校音乐专业学生 · 青年泛音乐受众",
        "",
        "【差异化优势】",
        "  ● 主打热门流行曲、影视金曲的民族乐器室内乐改编",
        "  ● 兼顾国风韵味、演奏质感与网络传播性",
        "  ● 短视频高光片段引流，适配国内各平台用户观看习惯",
        "",
        "【平台选择】",
        "依托抖音、视频号等国内平台进行本土化传播，打破传统民乐受众小众的局限，打造年轻化、专业化的国风音乐内容账号。",
    ]
)

# 2-2 运营过程
content_slide(
    "第二章 · 账号运营与传播成效",
    "二、运营过程与内容发布节奏",
    [
        "【三阶段运营模型】",
        "",
        "① 冷启动期",
        "   完成账号搭建与首批改编作品发布，借助国风、民乐话题标签积累初始粉丝",
        "",
        "② 成长期",
        "   根据平台数据优化选题，收集观众点歌需求，联动国内同领域账号扩大曝光",
        "",
        "③ 稳定运营期",
        "   固定更新频次，推出系列改编内容，维护粉丝社群，沉淀核心受众",
        "",
        "【发布节奏】",
        "  ● 发布时段：晚间 19:00–21:00（国内用户休闲高峰）",
        "  ● 更新频率：每周 2–3 条短视频",
        "  ● 内容流程：曲目选题 → 民乐室内乐改编 → 录制拍摄 → 剪辑适配 → 数据复盘",
        "  ● 及时回复评论互动，遵守国内平台版权规范，提升账号粘性与流量权重",
    ]
)

# 2-3 核心数据
s = prs.slides.add_slide(BLANK)
bg(s, C_MID)
add_rect(s, 0, 0, 13.33, 1.0, C_DARK)
add_rect(s, 0, 0, 0.08, 7.5, C_ACCENT)
accent_line(s, 0.08, 1.0, 13.25, 0.04)
add_textbox(s, "第二章 · 账号运营与传播成效", 0.2, 0.1, 8, 0.45, font_size=13, color=C_ACCENT)
add_textbox(s, "三、核心数据表现与传播效果分析", 0.2, 0.5, 12.5, 0.45,
            font_size=20, bold=True, color=C_WHITE)

# 数据卡片
cards = [
    ("17万+", "累计播放量"),
    ("1000+", "精准粉丝"),
    ("18-35岁", "核心受众年龄段"),
    ("多省市", "受众地域覆盖"),
]
for i, (val, label) in enumerate(cards):
    x = 0.4 + i * 3.2
    add_rect(s, x, 1.2, 2.8, 1.6, RGBColor(0x0D,0x1A,0x35))
    accent_line(s, x, 1.2, 2.8, 0.05)
    add_textbox(s, val,   x, 1.35, 2.8, 0.8,
                font_size=32, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)
    add_textbox(s, label, x, 2.15, 2.8, 0.45,
                font_size=14, color=C_LGRAY, align=PP_ALIGN.CENTER)

add_para_text(s, [
    "传播成效亮点：",
    "  ● 运营期间在国内各平台累计播放量达 17 万余次，精准粉丝突破千人",
    "  ● 多条作品进入国风、民乐垂类流量池",
    "  ● 突破传统民乐专业圈层限制，实现民族音乐大众化网络传播",
    "  ● 作品被大量转发至社交社群、校园群组，形成自发二次传播",
    "",
    "现存问题：",
    "  ● 视频开头吸引力不足    ● 内容产能较低",
    "  ● 传播能力有待提升      ● 团队协作仍有磨合空间",
], 0.35, 2.95, 12.6, 4.3, font_size=13)

# 2-4 受众反馈
content_slide(
    "第二章 · 账号运营与传播成效",
    "四、受众情况与用户反馈总结",
    [
        "【受众画像】",
        "  ● 年龄层：以国内 18–35 岁青年及高校学生为主",
        "  ● 兴趣特征：国风爱好者、民乐学习者、泛音乐受众",
        "  ● 用户认可度高，互动意愿强",
        "",
        "【用户反馈】",
        "  ● 评论区对民乐改编形式给予较高认可，不乏对后续创作的期待与点歌需求",
        "  ● 调查问卷显示用户对民族风流行音乐融合内容具有持续消费意愿",
        "  ● 部分用户主动将视频分享至社交社群，形成口碑传播",
        "",
        "【传播覆盖】",
        "  ● 覆盖国内多省市受众，突破传统民乐专业圈层限制",
        "  ● 作品被转发至校园群组，助力传统民乐在青年群体中的年轻化传承",
        "  ● 提升了 GongfuMusic 在国内国风音乐圈的影响力",
    ]
)


# ═══════════════════════════════════════════════════════════
# 7. 第三章
# ═══════════════════════════════════════════════════════════
chapter_slide("第三章", "民族风流行音乐改编实践")

# 3-1 作品完成情况
s = prs.slides.add_slide(BLANK)
bg(s, C_MID)
add_rect(s, 0, 0, 13.33, 1.0, C_DARK)
add_rect(s, 0, 0, 0.08, 7.5, C_ACCENT)
accent_line(s, 0.08, 1.0, 13.25, 0.04)
add_textbox(s, "第三章 · 民族风流行音乐改编实践", 0.2, 0.1, 8, 0.45, font_size=13, color=C_ACCENT)
add_textbox(s, "一、作品完成情况与选曲思路", 0.2, 0.5, 12.5, 0.45,
            font_size=20, bold=True, color=C_WHITE)

instrument_groups = [
    ("竹  笛", C_ACCENT,
     ["《青鸟》", "《Wind》", "《Rest Assured》", "《Anxiety》", "《Leave The Door Open》"]),
    ("唢  呐", RGBColor(0x7E,0xC8,0xA0),
     ["《柯南系列主题曲》", "《Bad Romance》", "《Somebody To Someone》"]),
    ("马头琴", RGBColor(0xA0,0xC4,0xFF),
     ["《Experience》", "《我爱你》", "《人生的旋转木马》", "《Salvatore》", "《爱琴海》"]),
]
for i, (inst, color, songs) in enumerate(instrument_groups):
    x = 0.35 + i * 4.3
    add_rect(s, x, 1.2, 4.0, 0.5, RGBColor(0x0D,0x1A,0x35))
    accent_line(s, x, 1.2, 4.0, 0.05)
    add_textbox(s, inst, x, 1.22, 4.0, 0.4,
                font_size=18, bold=True, color=color, align=PP_ALIGN.CENTER)
    for j, song in enumerate(songs):
        add_textbox(s, "♩ " + song, x+0.1, 1.85 + j*0.85, 3.8, 0.55,
                    font_size=14, color=C_LGRAY)

add_textbox(s, "选曲思路：优先选择大众认可度高的流行与国风曲目，兼顾民族乐器音色特质与旋律适配性，在保留原曲熟悉度的基础上融入民族音乐的旋律、音色与文化特色。",
            0.35, 6.3, 12.6, 0.9, font_size=12.5, color=RGBColor(0xA0,0xB8,0xD0))

# 3-2 乐器改编与改编特色
content_slide(
    "第三章 · 民族风流行音乐改编实践",
    "二、乐器改编与改编特色——以《爱情海》为例",
    [
        "以周杰伦经典曲目《爱情海》为原创基底，聚焦传统戏曲年轻化创新与民族乐器改编重构，打造流行音乐与越剧艺术跨界融合的全新版本。",
        "",
        "【改编方式】",
        "① 旋律改编：保留《爱情海》抒情细腻的情感内核与经典旋律线条，将越剧婉转柔美的唱腔韵律、板式节奏融入编曲，注入越剧独有的"水袖般"的东方韵律",
        "",
        "② 乐器改编：摒弃传统流行乐电子配器，以民乐竹笛、唢呐等江南丝竹乐器为核心，用民族器乐重构和声、织体与伴奏层次，让流行旋律扎根于传统戏曲器乐肌理",
        "",
        "【改编特色】",
        "以青年视角打破流行音乐与传统戏曲的壁垒，用"流行经典+民族器乐改编"的模式，让大众熟悉的旋律搭载非遗戏曲之美，既保留传唱度，又赋予作品浓郁的东方美学与文化辨识度。",
        "契合青苗计划鼓励青年创新、文化创业的宗旨，探索传统戏曲年轻化传播与原创音乐跨界改编的可落地路径。",
    ]
)

# 3-3 作品呈现与录制
content_slide(
    "第三章 · 民族风流行音乐改编实践",
    "三、作品呈现与录制过程",
    [
        "【录制流程】",
        "  曲目确定 → 编曲写谱 → 分轨录音 → 混音后期 → 视频拍摄 → 剪辑发布",
        "",
        "【乐器录制特点】",
        "  ● 竹笛：注重气息控制与音色圆润度，捕捉民乐清透音色",
        "  ● 唢呐：突出穿透力与戏曲韵味，兼顾流行融合感",
        "  ● 马头琴：展现草原风情与深沉情感，强调弓法表现力",
        "",
        "【视频呈现】",
        "  ● 以演奏实录为主，简洁构图凸显民族乐器质感",
        "  ● 配合字幕与后期调色，营造国风视觉氛围",
        "  ● 适配短视频传播节奏，前 3 秒设置高光吸睛内容",
        "",
        "【代表作品】",
        "  竹笛版《珠玉》《魔女宅急便》等，兼具专业演奏质感与大众审美接受度",
    ]
)

# 3-4 传播效果与评价
content_slide(
    "第三章 · 民族风流行音乐改编实践",
    "四、作品传播效果与评价",
    [
        "【数据表现】",
        "  ● 已发布作品累计播放量 17 万余次",
        "  ● 多条作品进入垂类流量池，实现精准受众触达",
        "  ● 用户评论中对民乐改编形式认可度高，互动积极",
        "",
        "【用户评价】",
        "  ● 对民族乐器音色与流行旋律融合的新鲜感给予正面反馈",
        "  ● 不乏专业乐手与音乐爱好者的认可与转发",
        "  ● 部分用户留言点歌，形成内容需求反馈闭环",
        "",
        "【文化价值评价】",
        "  ● 验证了"流行旋律+民族乐器"低门槛传播模式的可行性",
        "  ● 有效降低大众接触民族音乐的心理门槛",
        "  ● 为民族音乐数字化传播提供了可复制的轻量化实践样本",
        "",
        "【创作困难与解决方案】",
        "  ● 困难：民族乐器音域与流行编曲的适配平衡 → 解决：反复试奏调整音区与配器层次",
        "  ● 困难：录音设备条件限制 → 解决：优化录音环境与后期混音处理",
    ]
)


# ═══════════════════════════════════════════════════════════
# 8. 第四章
# ═══════════════════════════════════════════════════════════
chapter_slide("第四章", "项目整体总结与发展展望")

# 4-1 项目总结
content_slide(
    "第四章 · 项目整体总结与发展展望",
    "一、项目整体总结（上）",
    [
        "本项目以"铸牢中华民族共同体意识：中华民族多民族音乐数字化创新传播路径研究与实践"为核心主题，由中央民族大学音乐学院主导，历经前期筹备、中期实施与后期总结三个阶段，圆满完成各项既定工作任务。",
        "",
        "【项目背景与定位】",
        "确立"民族乐器+流行音乐+数字化"创新传播路径，以"KungFuMusic" IP 为核心，明确学术价值（丰富中华民族共同体意识的教育实践）与应用价值（推动民族乐器普及）。",
        "",
        "【团队建设】",
        "  ● 涵盖作曲理论、音乐表演、计算机等多专业核心团队",
        "  ● 成员具备演奏、作曲编曲、新媒体运营、视频拍摄、撰稿等综合能力",
        "  ● 牛小凤老师担任指导老师，提供全流程专业指导",
        "",
        "【实施流程】",
        "严格遵循"筛选—改编—传播—评估"全流程模式，构建"线上传播+线下联动"传播格局。",
    ]
)

content_slide(
    "第四章 · 项目整体总结与发展展望",
    "一、项目整体总结（下）",
    [
        "【阶段性成果】",
        "  ● 民乐演奏视频累计获得 17w+ 用户观看，精准粉丝 1000+",
        "  ● 成功搭建"KungFuMusic"核心抖音账号，形成完整的民族风流行音乐室内乐转化与传播体系",
        "  ● 产出竹笛版《珠玉》《魔女宅急便》等优质作品",
        "  ● 探索出适配国内传播的民族音乐创新路径",
        "  ● 团队成员专业能力与协作能力显著提升",
        "",
        "【主要不足】",
        "  ● 全网粉丝规模未达预期，内容转化率仍有提升空间",
        "  ● 部分作品未能充分兼顾民族内核与流行听感的平衡",
        "  ● 传播渠道多元化程度不足，对不同平台的策略优化不够精准",
        "",
        "总体而言，本项目不仅探索了民族音乐数字化创新传播的有效路径，也为民族音乐专业学生搭建了专业实践平台，具有重要的现实意义与学术价值。",
    ]
)

# 4-2 不足与反思
content_slide(
    "第四章 · 项目整体总结与发展展望",
    "二、项目不足与反思",
    [
        "① 目标设定与实际成果存在差距",
        "   项目前期立项目标较为远大，受实践周期、团队合作及制作条件限制，实际完成成果与最初理想目标有一定差距。",
        "",
        "② 作品改编创作层面仍有提升空间",
        "   改编题材与乐器搭配类型不够丰富，局限于竹笛、唢呐、马头琴；民族音乐素材挖掘不够深入；作品打磨精细度和内容质量需进一步提升。",
        "",
        "③ 账号运营管理不够系统规范",
        "   内容形式较为单一，多以纯演奏视频为主，缺少乐器科普、改编思路解读、创作幕后等衍生内容；日常粉丝维护、话题策划、线上互动活动不足；更新节奏和运营精细化程度有待加强。",
        "",
        "④ 传播渠道较为单一",
        "   主要依赖抖音单个短视频平台，未搭建多平台传播矩阵；对外联动音乐院校师生、民乐爱好者投稿合作的力度不足，未能形成规模化内容生产与大范围传播效应。",
        "",
        "⑤ 跨专业团队协作仍需磨合",
        "   各环节沟通效率、作品产出效率以及成品精细化制作水平，都有待在实践中继续优化完善。",
    ]
)

# 4-3 未来展望
content_slide(
    "第四章 · 项目整体总结与发展展望",
    "三、未来展望",
    [
        "【内容创作】",
        "  ● 每月坚持完成 2-3 首民族乐器改编曲目",
        "  ● 产出隔天更新的优质短视频",
        "  ● 新增民乐科普类轻量化内容，持续优化"KungFuMusic"抖音账号运营",
        "",
        "【传播推广】",
        "  ● 在深耕抖音的基础上，拓展视频号、小红书两大平台",
        "  ● 定制差异化传播内容，完善学生投稿激励机制",
        "  ● 吸引校内专业学生参与，壮大传播队伍",
        "",
        "【团队建设】",
        "  ● 每季度开展复盘交流，邀请指导老师针对性指导",
        "  ● 深化跨专业协作，提升成员专业能力",
        "",
        "我们将持续推动民族音乐普及，切实发挥项目实践价值，为助力中华优秀传统文化传承与发展，贡献青年力量。",
    ]
)


# ═══════════════════════════════════════════════════════════
# 9. 结束页
# ═══════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, C_DARK)
add_rect(s, 0, 0, 0.08, 7.5, C_ACCENT)
add_rect(s, 0.08, 3.1, 13.25, 1.3, RGBColor(0x0D, 0x14, 0x26))
accent_line(s, 0.08, 3.1,  13.25, 0.06)
accent_line(s, 0.08, 4.4,  13.25, 0.06)
add_textbox(s, "感谢聆听", 0, 3.2, 13.33, 0.9,
            font_size=44, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
add_textbox(s, "KungfuMusic · 中央民族大学音乐学院 · 青苗计划", 0, 4.55, 13.33, 0.5,
            font_size=17, color=C_SUB, align=PP_ALIGN.CENTER)
add_textbox(s, "让民族乐器之声融入当代文化传播场景", 0, 5.15, 13.33, 0.5,
            font_size=14, color=RGBColor(0x70,0x80,0x90), align=PP_ALIGN.CENTER)


# ── 保存 ─────────────────────────────────────────────────
out = "d:/Users/vs code 代码/青苗计划PPT.pptx"
prs.save(out)
print("saved:", out)
