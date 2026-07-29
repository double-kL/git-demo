# -*- coding: utf-8 -*-
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree
import json

# ── Palette ──────────────────────────────────────────────────────────────────
def rgb(r, g, b): return RGBColor(r, g, b)

DARK1  = (6,  10, 24)    # Gradient top
DARK2  = (13, 21, 50)    # Gradient bottom
CARD   = rgb(17, 26, 58) # Content card
HDR    = rgb( 4,  7, 17) # Header / footer band
GOLD   = rgb(212,175, 55)
GOLD2  = rgb(148,116, 28)
ICE    = rgb(155,205,242)
WWHITE = rgb(242,238,230)
GRAY   = rgb(145,160,178)
DIM    = rgb( 52, 72,105)
WMARK  = rgb( 18, 32, 72)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

# ── Primitives ────────────────────────────────────────────────────────────────
def lerp(a, b, t): return int(a + (b - a) * t)

def grad_bg(slide, top=DARK1, bot=DARK2, steps=16):
    h = 7.5 / steps
    for i in range(steps):
        t = i / max(steps - 1, 1)
        c = rgb(lerp(top[0],bot[0],t), lerp(top[1],bot[1],t), lerp(top[2],bot[2],t))
        sh = slide.shapes.add_shape(1, Inches(0), Inches(i*h), Inches(13.33), Inches(h+0.012))
        sh.line.fill.background()
        sh.fill.solid()
        sh.fill.fore_color.rgb = c

def box(slide, l, t, w, h, color):
    sh = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    sh.line.fill.background()
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    return sh

def text(slide, s, l, t, w, h, size=16, bold=False,
         color=WWHITE, align=PP_ALIGN.LEFT, italic=False):
    b = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    b.word_wrap = True
    tf = b.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = s
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    r.font.name = "Microsoft YaHei"
    return b

def body_text(slide, lines, l, t, w, h, size=13, sp=24):
    b = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    b.word_wrap = True
    tf = b.text_frame
    tf.word_wrap = True
    first = True
    for line in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = PP_ALIGN.LEFT
        pr = p._p.get_or_add_pPr()
        ls = etree.SubElement(pr, qn("a:lnSpc"))
        sv = etree.SubElement(ls, qn("a:spcPts"))
        sv.set("val", str(sp * 100))
        r = p.add_run()
        r.text = line
        is_h = line.startswith("【") and "】" in line
        r.font.size = Pt(size + (0.5 if is_h else 0))
        r.font.bold = is_h
        r.font.color.rgb = WWHITE if is_h else GRAY
        r.font.name = "Microsoft YaHei"
    return b

def gold_rule(slide, l, t, w, h=0.038):
    box(slide, l, t, w, h, GOLD)

def dim_rule(slide, l, t, w):
    box(slide, l, t, w, 0.016, DIM)

def footer(slide, label):
    box(slide, 0.06, 7.28, 13.27, 0.22, HDR)
    box(slide, 0.06, 7.28, 13.27, 0.016, DIM)
    text(slide, label, 0.22, 7.3, 10, 0.18, size=9, color=GOLD2, italic=True)

# ── COVER ─────────────────────────────────────────────────────────────────────
def make_cover(data):
    s = prs.slides.add_slide(BLANK)
    grad_bg(s, (4, 7, 18), (11, 17, 44))

    # Full-height left gold bar + shadow bar
    box(s, 0,     0, 0.06, 7.5, GOLD)
    box(s, 0.06,  0, 0.025, 7.5, GOLD2)

    # Top thin gold line
    gold_rule(s, 0, 0, 13.33, 0.05)

    # Bottom band
    box(s, 0, 6.62, 13.33, 0.88, HDR)
    gold_rule(s, 0, 6.62, 13.33, 0.04)

    # Right-side decorative column: five thin bars
    for i, (x, ht) in enumerate([(12.5,4.2),(12.62,3.6),(12.74,5.0),(12.86,3.0),(12.98,4.5)]):
        box(s, x, (7.5-ht)/2, 0.07, ht, rgb(lerp(18,40,i/4), lerp(28,60,i/4), lerp(58,100,i/4)))

    # Top tag
    text(s, "KungfuMusic", 1.25, 1.35, 8, 0.42,
         size=15, color=GOLD2, italic=True)

    # Main title
    text(s, data["cover"]["title"], 1.25, 1.82, 10.5, 1.65,
         size=72, bold=True, color=GOLD, align=PP_ALIGN.LEFT)

    # Gold rule under title
    gold_rule(s, 1.25, 3.5, 9.8, 0.042)

    # Subtitle
    text(s, "民族风流行音乐室内乐改编与数字化传播实践",
         1.25, 3.62, 10.5, 0.65, size=22, color=ICE)

    # Key words
    text(s, "民族乐器  ·  室内乐改编  ·  抖音短视频  ·  文化传播",
         1.25, 4.38, 10, 0.42, size=13, color=GRAY, italic=True)

    # Bottom: org + year
    text(s, data["cover"]["org"], 1.25, 6.72, 8, 0.38, size=13.5, color=GRAY)
    text(s, data["cover"]["year"], 11.5, 6.72, 1.7, 0.38,
         size=13.5, color=GRAY, align=PP_ALIGN.RIGHT)

# ── CHAPTER TRANSITION ───────────────────────────────────────────────────────
def make_chapter(title, subtitle, ch_num=None):
    s = prs.slides.add_slide(BLANK)
    grad_bg(s, (3, 6, 16), (10, 16, 40))

    # Watermark chapter number (right side, barely visible)
    if ch_num:
        text(s, "0" + str(ch_num), 6.5, -0.5, 7.0, 8.5,
             size=200, bold=True, color=WMARK, align=PP_ALIGN.LEFT)

    # Left gold bar
    box(s, 0, 0, 0.06, 7.5, GOLD)
    box(s, 0.06, 0, 0.025, 7.5, GOLD2)

    # Top thin line
    gold_rule(s, 0.06, 0, 13.27, 0.05)

    # Two horizontal gold rules framing title
    gold_rule(s, 1.4, 3.0, 10.8, 0.04)
    gold_rule(s, 1.4, 4.7, 10.8, 0.04)

    # Chapter label (small)
    if ch_num:
        label = "CHAPTER  0" + str(ch_num)
    else:
        label = "SECTION"
    text(s, label, 1.4, 2.55, 6, 0.38,
         size=11, color=GOLD2, italic=True)

    # Main title
    text(s, title, 1.4, 3.08, 10.8, 1.52,
         size=44, bold=True, color=WWHITE, align=PP_ALIGN.LEFT)

    # Subtitle
    if subtitle:
        text(s, subtitle, 1.4, 4.78, 10.8, 0.72,
             size=21, color=ICE, align=PP_ALIGN.LEFT)

    return s

# ── CONTENT SLIDE ────────────────────────────────────────────────────────────
def make_content(section, heading, paras, fs=13.5):
    s = prs.slides.add_slide(BLANK)
    grad_bg(s, DARK1, DARK2)

    # Header band
    box(s, 0, 0, 13.33, 0.98, HDR)
    gold_rule(s, 0, 0.94, 13.33, 0.04)

    # Left accent bar
    box(s, 0, 0, 0.06, 7.5, GOLD)
    box(s, 0.06, 0, 0.025, 7.5, GOLD2)

    # Breadcrumb
    text(s, section, 0.2, 0.12, 11, 0.36,
         size=10.5, color=GOLD2, italic=True)

    # Heading
    text(s, heading, 0.2, 0.5, 12.5, 0.44,
         size=21, bold=True, color=WWHITE)

    # Content card
    box(s, 0.14, 1.08, 13.05, 6.1, CARD)
    box(s, 0.14, 1.08, 13.05, 0.045, GOLD)  # gold top edge on card

    # Body text
    body_text(s, paras, 0.38, 1.22, 12.6, 5.85, size=fs, sp=23)

    footer(s, section)
    return s

# ── TOC SLIDE ────────────────────────────────────────────────────────────────
def make_toc(entries):
    s = prs.slides.add_slide(BLANK)
    grad_bg(s, DARK1, DARK2)
    box(s, 0, 0, 13.33, 0.98, HDR)
    gold_rule(s, 0, 0.94, 13.33, 0.04)
    box(s, 0, 0, 0.06, 7.5, GOLD)
    box(s, 0.06, 0, 0.025, 7.5, GOLD2)

    text(s, "目  录 / Contents", 0.2, 0.5, 12, 0.44,
         size=21, bold=True, color=WWHITE)

    c_list = [GOLD, ICE, rgb(150,220,155), rgb(220,165,115)]

    for i, e in enumerate(entries):
        y = 1.1 + i * 1.52
        c = c_list[i % len(c_list)]

        # Number badge
        box(s, 0.18, y + 0.04, 0.62, 0.62, c)
        text(s, str(i+1), 0.18, y + 0.04, 0.62, 0.62,
             size=24, bold=True, color=HDR, align=PP_ALIGN.CENTER)

        # Chapter tag
        text(s, e["ch"], 0.98, y + 0.06, 2.0, 0.34,
             size=12, bold=True, color=c)

        # Title
        text(s, e["title"], 0.98, y + 0.36, 5.8, 0.48,
             size=18, bold=True, color=WWHITE)

        # Sub keywords
        text(s, e["sub"], 7.1, y + 0.15, 6.05, 0.52,
             size=11.5, color=GRAY, italic=True)

        # Separator
        if i < len(entries) - 1:
            dim_rule(s, 0.18, y + 0.88, 12.95)

    footer(s, "目录")
    return s

# ── DATA CARDS SLIDE ─────────────────────────────────────────────────────────
def make_data_cards(section, heading, cards, paras):
    s = prs.slides.add_slide(BLANK)
    grad_bg(s, DARK1, DARK2)
    box(s, 0, 0, 13.33, 0.98, HDR)
    gold_rule(s, 0, 0.94, 13.33, 0.04)
    box(s, 0, 0, 0.06, 7.5, GOLD)
    box(s, 0.06, 0, 0.025, 7.5, GOLD2)
    text(s, section, 0.2, 0.12, 11, 0.36, size=10.5, color=GOLD2, italic=True)
    text(s, heading, 0.2, 0.5, 12.5, 0.44, size=21, bold=True, color=WWHITE)

    c_list = [GOLD, ICE, rgb(150,218,150), rgb(220,162,110)]

    for i, card in enumerate(cards):
        c = c_list[i % len(c_list)]
        x = 0.18 + i * 3.24

        # Card
        box(s, x, 1.08, 3.05, 1.72, CARD)
        box(s, x, 1.08, 3.05, 0.055, c)   # color top edge

        # Large value
        text(s, card["value"], x, 1.18, 3.05, 0.9,
             size=36, bold=True, color=c, align=PP_ALIGN.CENTER)

        # Label
        text(s, card["label"], x, 2.08, 3.05, 0.52,
             size=13, color=GRAY, align=PP_ALIGN.CENTER)

    # Content
    box(s, 0.14, 2.95, 13.05, 4.22, CARD)
    box(s, 0.14, 2.95, 13.05, 0.04, DIM)
    body_text(s, paras, 0.38, 3.08, 12.6, 4.0, size=13, sp=22)

    footer(s, section)
    return s

# ── INSTRUMENT SLIDE ─────────────────────────────────────────────────────────
def make_instrument(section, heading, groups, note):
    s = prs.slides.add_slide(BLANK)
    grad_bg(s, DARK1, DARK2)
    box(s, 0, 0, 13.33, 0.98, HDR)
    gold_rule(s, 0, 0.94, 13.33, 0.04)
    box(s, 0, 0, 0.06, 7.5, GOLD)
    box(s, 0.06, 0, 0.025, 7.5, GOLD2)
    text(s, section, 0.2, 0.12, 11, 0.36, size=10.5, color=GOLD2, italic=True)
    text(s, heading, 0.2, 0.5, 12.5, 0.44, size=21, bold=True, color=WWHITE)

    c_list = [GOLD, ICE, rgb(150,218,150)]

    for i, grp in enumerate(groups):
        c = c_list[i % len(c_list)]
        x = 0.18 + i * 4.38

        # Group card
        box(s, x, 1.08, 4.1, 5.15, CARD)
        box(s, x, 1.08, 4.1, 0.055, c)

        # Name
        text(s, grp["name"], x, 1.18, 4.1, 0.6,
             size=21, bold=True, color=c, align=PP_ALIGN.CENTER)

        # Divider
        dim_rule(s, x + 0.15, 1.88, 3.8)

        # Songs
        for j, song in enumerate(grp["songs"]):
            text(s, song, x + 0.2, 2.02 + j * 0.82, 3.7, 0.62,
                 size=13.5, color=GRAY)

    # Note
    if note:
        text(s, note, 0.2, 6.32, 12.9, 0.72,
             size=11, color=DIM, italic=True)

    footer(s, section)
    return s

# ── TWO-COL SLIDE ────────────────────────────────────────────────────────────
def make_two_col(section, heading, lt, rt, left, right):
    s = prs.slides.add_slide(BLANK)
    grad_bg(s, DARK1, DARK2)
    box(s, 0, 0, 13.33, 0.98, HDR)
    gold_rule(s, 0, 0.94, 13.33, 0.04)
    box(s, 0, 0, 0.06, 7.5, GOLD)
    box(s, 0.06, 0, 0.025, 7.5, GOLD2)
    text(s, section, 0.2, 0.12, 11, 0.36, size=10.5, color=GOLD2, italic=True)
    text(s, heading, 0.2, 0.5, 12.5, 0.44, size=21, bold=True, color=WWHITE)

    # Left panel
    box(s, 0.14, 1.08, 6.35, 6.1, CARD)
    box(s, 0.14, 1.08, 6.35, 0.05, GOLD)

    # Right panel
    box(s, 6.64, 1.08, 6.55, 6.1, CARD)
    box(s, 6.64, 1.08, 6.55, 0.05, ICE)

    if lt:
        text(s, lt, 0.32, 1.18, 6.0, 0.44, size=14, bold=True, color=GOLD)
    if rt:
        text(s, rt, 6.82, 1.18, 6.1, 0.44, size=14, bold=True, color=ICE)

    top = 1.7 if lt else 1.2
    body_text(s, left,  0.32, top, 6.0, 6.5 - top + 1.08, size=13, sp=22)
    body_text(s, right, 6.82, top, 6.1, 6.5 - top + 1.08, size=13, sp=22)

    footer(s, section)
    return s

# ── END SLIDE ────────────────────────────────────────────────────────────────
def make_end(data):
    s = prs.slides.add_slide(BLANK)
    grad_bg(s, (4, 7, 18), (11, 17, 44))

    gold_rule(s, 0, 0, 13.33, 0.05)

    box(s, 0, 0, 0.06, 7.5, GOLD)
    box(s, 0.06, 0, 0.025, 7.5, GOLD2)

    # Right decorative bars (mirror of cover)
    for i, (x, ht) in enumerate([(12.5,4.2),(12.62,3.6),(12.74,5.0),(12.86,3.0),(12.98,4.5)]):
        box(s, x, (7.5-ht)/2, 0.07, ht, rgb(lerp(18,40,i/4), lerp(28,60,i/4), lerp(58,100,i/4)))

    # Decorative dots above title
    for i in range(7):
        box(s, 4.8 + i * 0.52, 2.62, 0.14, 0.14, GOLD2 if i % 2 else GOLD)

    # Two gold rules framing title
    gold_rule(s, 1.5, 2.88, 10.5, 0.042)
    gold_rule(s, 1.5, 4.58, 10.5, 0.042)

    # Title
    text(s, data["end"]["title"], 0, 3.0, 13.33, 1.42,
         size=52, bold=True, color=WWHITE, align=PP_ALIGN.CENTER)

    # Sub 1
    text(s, data["end"]["sub1"], 0, 4.68, 13.33, 0.52,
         size=16, color=ICE, align=PP_ALIGN.CENTER)

    # Sub 2
    text(s, data["end"]["sub2"], 0, 5.3, 13.33, 0.45,
         size=13, color=GRAY, italic=True, align=PP_ALIGN.CENTER)

    # Bottom band
    box(s, 0.06, 6.88, 13.27, 0.62, HDR)
    gold_rule(s, 0.06, 6.88, 13.27, 0.04)
    text(s, "KungfuMusic  ·  中央民族大学音乐学院  ·  青苗计划",
         0, 7.0, 13.33, 0.38, size=12, color=GRAY, align=PP_ALIGN.CENTER)

# ─── MAIN ────────────────────────────────────────────────────────────────────
with open("slides_content.json", "r", encoding="utf-8") as f:
    data = json.load(f)

# Cover
make_cover(data)

# Preface + Abstract
for ch in data["chapters"]:
    make_chapter(ch["chapter_title"], ch.get("subtitle", ""), ch_num=None)
    for sl in ch["slides"]:
        make_content(ch["section_label"], sl["heading"], sl["body"])

# TOC
make_chapter(data["toc"]["chapter_title"], data["toc"].get("subtitle", ""), ch_num=None)
make_toc(data["toc"]["entries"])

# Main chapters
for i, ch_data in enumerate(data["main_chapters"], 1):
    make_chapter(ch_data["chapter_title"], ch_data.get("subtitle", ""), ch_num=i)
    for sl in ch_data["slides"]:
        if sl.get("two_col"):
            make_two_col(ch_data["section_label"], sl["heading"],
                         sl.get("left_title",""), sl.get("right_title",""),
                         sl["left"], sl["right"])
        elif sl.get("data_cards"):
            make_data_cards(ch_data["section_label"], sl["heading"],
                            sl["cards"], sl["body"])
        elif sl.get("instrument_groups"):
            make_instrument(ch_data["section_label"], sl["heading"],
                            sl["groups"], sl.get("note",""))
        else:
            make_content(ch_data["section_label"], sl["heading"], sl["body"])

# End
make_end(data)

out = "d:/Users/vs code 代码/青苗计划PPT.pptx"
prs.save(out)
print("saved:", out, "  slides:", len(prs.slides))
